#!/usr/bin/env python3
"""
THALIA - Creative Fleet Designer Agent
Port: 8032

Visual design, layouts, presentations, and web/UI design for LeverEdge.
Part of the Creative Fleet alongside CALLIOPE (Writer), ERATO (Multimedia),
CLIO (Research Writer), and MUSE (Orchestrator).

CAPABILITIES:
- Presentation design (PowerPoint/slides)
- Document formatting (Word/PDF)
- Chart/graph creation
- Brand template application
- Layout optimization
- Thumbnail design
- Landing page design (HTML/CSS/Tailwind)
- UI component design (buttons, forms, cards, navbars)
- Website wireframes/mockups
- Responsive website templates

NOT primarily LLM-powered - focuses on programmatic generation.
"""

import os
import sys
import json
import uuid
import httpx
import tempfile
import base64
from pathlib import Path
from datetime import datetime, date
from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import Optional, Dict, Any, List, Literal, Union
from io import BytesIO

# Image and document libraries
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

# Add shared modules to path
sys.path.append('/opt/leveredge/control-plane/shared')
try:
    from cost_tracker import CostTracker
except ImportError:
    # Fallback if shared module not available
    class CostTracker:
        def __init__(self, name): self.name = name
        async def log_usage(self, **kwargs): pass

app = FastAPI(
    title="THALIA",
    description="Creative Fleet Designer - Visual design, layouts, presentations",
    version="1.0.0"
)

# =============================================================================
# CONFIGURATION
# =============================================================================

EVENT_BUS_URL = os.getenv("EVENT_BUS_URL", "http://event-bus:8099")
OUTPUT_DIR = os.getenv("THALIA_OUTPUT_DIR", "/tmp/thalia_output")
ASSETS_DIR = os.getenv("THALIA_ASSETS_DIR", "/opt/leveredge/assets")

# Ensure output directory exists
Path(OUTPUT_DIR).mkdir(parents=True, exist_ok=True)

# Agent endpoints
AGENT_ENDPOINTS = {
    "MUSE": "http://muse:8030",
    "CALLIOPE": "http://calliope:8031",
    "ERATO": "http://erato:8033",
    "CLIO": "http://clio:8034",
    "HEPHAESTUS": "http://hephaestus:8011",
    "EVENT_BUS": EVENT_BUS_URL
}

# Key dates
LAUNCH_DATE = date(2026, 3, 1)

# Initialize cost tracker
cost_tracker = CostTracker("THALIA")

# =============================================================================
# DESIGN SYSTEM
# =============================================================================

COLOR_PALETTES = {
    "leveredge": {
        "primary": "#1B2951",      # Deep navy blue
        "secondary": "#B8860B",    # Dark goldenrod
        "accent": "#36454F",       # Charcoal
        "background": "#FFFFFF",   # White
        "text": "#1A1A1A",         # Near black
        "light": "#F5F7FA",        # Light gray
        "success": "#28A745",      # Green
        "warning": "#FFC107",      # Amber
        "error": "#DC3545",        # Red
    },
    "professional": {
        "primary": "#2C3E50",
        "secondary": "#3498DB",
        "accent": "#E74C3C",
        "background": "#FFFFFF",
        "text": "#2C3E50",
        "light": "#ECF0F1",
        "success": "#27AE60",
        "warning": "#F39C12",
        "error": "#E74C3C",
    },
    "modern": {
        "primary": "#6C63FF",
        "secondary": "#FF6584",
        "accent": "#43D9AD",
        "background": "#FFFFFF",
        "text": "#2D3748",
        "light": "#F7FAFC",
        "success": "#48BB78",
        "warning": "#ECC94B",
        "error": "#F56565",
    },
    "minimal": {
        "primary": "#000000",
        "secondary": "#666666",
        "accent": "#FF4444",
        "background": "#FFFFFF",
        "text": "#333333",
        "light": "#F9F9F9",
        "success": "#00C853",
        "warning": "#FFD600",
        "error": "#FF1744",
    },
    "dark": {
        "primary": "#BB86FC",
        "secondary": "#03DAC6",
        "accent": "#CF6679",
        "background": "#121212",
        "text": "#FFFFFF",
        "light": "#1E1E1E",
        "success": "#00E676",
        "warning": "#FFEA00",
        "error": "#FF5252",
    }
}

TYPOGRAPHY = {
    "headings": "Inter, Arial, sans-serif",
    "body": "Inter, Arial, sans-serif",
    "code": "JetBrains Mono, Consolas, monospace",
    "sizes": {
        "h1": 44,
        "h2": 32,
        "h3": 24,
        "h4": 20,
        "body": 16,
        "small": 14,
        "caption": 12
    }
}

SLIDE_LAYOUTS = {
    "title": "title_slide",
    "content": "content_slide",
    "two_column": "two_column",
    "image_left": "image_left",
    "image_right": "image_right",
    "quote": "quote_slide",
    "section": "section_header",
    "comparison": "comparison_slide",
    "timeline": "timeline_slide",
    "bullet_points": "bullet_points"
}

# =============================================================================
# TIME AWARENESS
# =============================================================================

def get_time_context() -> dict:
    """Get current time context for the agent"""
    now = datetime.now()
    today = now.date()
    days_to_launch = (LAUNCH_DATE - today).days

    return {
        "current_datetime": now.isoformat(),
        "current_date": today.isoformat(),
        "current_time": now.strftime("%I:%M %p"),
        "day_of_week": now.strftime("%A"),
        "launch_date": LAUNCH_DATE.isoformat(),
        "days_to_launch": days_to_launch,
        "launch_status": "LAUNCHED" if days_to_launch <= 0 else f"{days_to_launch} days to launch"
    }

# =============================================================================
# UTILITY FUNCTIONS
# =============================================================================

def hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def hex_to_pptx_rgb(hex_color: str) -> RgbColor:
    """Convert hex color to pptx RgbColor"""
    r, g, b = hex_to_rgb(hex_color)
    return RgbColor(r, g, b)

def hex_to_docx_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to docx RGBColor"""
    r, g, b = hex_to_rgb(hex_color)
    return RGBColor(r, g, b)

def get_palette(brand_id: Optional[str] = None) -> dict:
    """Get color palette by brand ID"""
    if brand_id and brand_id in COLOR_PALETTES:
        return COLOR_PALETTES[brand_id]
    return COLOR_PALETTES["leveredge"]

def generate_file_id() -> str:
    """Generate unique file ID"""
    return f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:8]}"

async def notify_event_bus(event_type: str, details: dict = {}):
    """Publish event to Event Bus"""
    try:
        time_ctx = get_time_context()
        async with httpx.AsyncClient() as client:
            await client.post(
                f"{EVENT_BUS_URL}/events",
                json={
                    "source_agent": "THALIA",
                    "event_type": event_type,
                    "details": details,
                    "timestamp": time_ctx["current_datetime"],
                    "days_to_launch": time_ctx["days_to_launch"]
                },
                timeout=5.0
            )
    except Exception as e:
        print(f"[THALIA] Event bus notification failed: {e}")

# =============================================================================
# PYDANTIC MODELS
# =============================================================================

class SlideContent(BaseModel):
    """Content for a single slide"""
    layout: str = "content"  # title, content, two_column, image_left, etc.
    title: Optional[str] = None
    subtitle: Optional[str] = None
    body: Optional[str] = None
    bullets: Optional[List[str]] = None
    image_path: Optional[str] = None
    notes: Optional[str] = None
    left_content: Optional[str] = None
    right_content: Optional[str] = None

class PresentationRequest(BaseModel):
    """Request to create a presentation"""
    content: Dict[str, Any]  # {slides: [...], title?, author?}
    brand_id: Optional[str] = "leveredge"
    template: Optional[str] = None
    style: Optional[str] = "professional"

class DocumentRequest(BaseModel):
    """Request to format a document"""
    content: str
    format: Literal["docx", "pdf"] = "docx"
    brand_id: Optional[str] = "leveredge"
    template: Optional[str] = None
    title: Optional[str] = None
    author: Optional[str] = None
    headers: Optional[Dict[str, str]] = None  # {header: "text", footer: "text"}

class ChartRequest(BaseModel):
    """Request to generate a chart"""
    type: Literal["bar", "line", "pie", "donut", "scatter", "area", "hbar"]
    data: Dict[str, Any]  # {labels: [...], values: [...], series?: [...]}
    style: Optional[str] = "leveredge"
    animated: Optional[bool] = False
    title: Optional[str] = None
    width: Optional[int] = 800
    height: Optional[int] = 600

class ThumbnailRequest(BaseModel):
    """Request to create a thumbnail"""
    title: str
    style: Optional[str] = "gradient"  # gradient, solid, pattern, image
    brand_id: Optional[str] = "leveredge"
    background_image: Optional[str] = None
    subtitle: Optional[str] = None
    icon: Optional[str] = None

class LayoutRequest(BaseModel):
    """Request for layout optimization"""
    content_type: Literal["email", "document", "slide", "social"]
    content: str
    brand_id: Optional[str] = "leveredge"
    constraints: Optional[Dict[str, Any]] = None

# =============================================================================
# WEB/UI DESIGN MODELS
# =============================================================================

class ContentSection(BaseModel):
    """A section of content for landing pages"""
    type: Literal["hero", "features", "cta", "testimonials", "pricing", "faq", "footer", "stats", "team", "contact", "custom"]
    title: Optional[str] = None
    subtitle: Optional[str] = None
    content: Optional[str] = None
    items: Optional[List[Dict[str, Any]]] = None  # For features, testimonials, etc.
    cta_text: Optional[str] = None
    cta_url: Optional[str] = "#"
    image_url: Optional[str] = None
    background: Optional[str] = None  # color or "gradient" or "image"

class LandingPageRequest(BaseModel):
    """Request to generate a landing page"""
    title: str
    sections: List[ContentSection]
    brand_id: Optional[str] = "leveredge"
    style: Optional[Literal["modern", "minimal", "corporate", "playful", "dark"]] = "modern"
    use_tailwind: Optional[bool] = True
    include_animations: Optional[bool] = True
    responsive: Optional[bool] = True
    custom_colors: Optional[Dict[str, str]] = None
    custom_fonts: Optional[Dict[str, str]] = None
    logo_url: Optional[str] = None
    favicon_url: Optional[str] = None

class UIComponentRequest(BaseModel):
    """Request to generate UI components"""
    component_type: Literal["button", "card", "form", "navbar", "footer", "modal", "alert", "badge", "input", "dropdown", "tabs", "accordion", "pricing_card", "testimonial_card", "feature_card", "hero", "cta_banner"]
    variant: Optional[str] = "primary"  # primary, secondary, outline, ghost, etc.
    size: Optional[Literal["sm", "md", "lg", "xl"]] = "md"
    content: Optional[Dict[str, Any]] = None  # Component-specific content
    brand_id: Optional[str] = "leveredge"
    use_tailwind: Optional[bool] = True
    include_javascript: Optional[bool] = False
    dark_mode: Optional[bool] = False

class WireframeRequest(BaseModel):
    """Request to generate a wireframe mockup"""
    page_type: Literal["landing", "dashboard", "blog", "ecommerce", "portfolio", "saas", "app"]
    sections: List[str]  # List of section names
    brand_id: Optional[str] = "leveredge"
    annotations: Optional[bool] = True
    format: Optional[Literal["html", "svg", "png"]] = "html"
    include_grid: Optional[bool] = True

class WebsiteTemplateRequest(BaseModel):
    """Request to generate a complete website template"""
    template_type: Literal["landing", "portfolio", "blog", "business", "saas", "agency", "startup"]
    pages: Optional[List[str]] = ["home"]  # home, about, services, contact, etc.
    brand_id: Optional[str] = "leveredge"
    style: Optional[Literal["modern", "minimal", "corporate", "playful", "dark"]] = "modern"
    use_tailwind: Optional[bool] = True
    include_components: Optional[List[str]] = None  # navbar, footer, etc.
    responsive_breakpoints: Optional[Dict[str, int]] = None
    custom_colors: Optional[Dict[str, str]] = None
    include_javascript: Optional[bool] = True

# =============================================================================
# PRESENTATION GENERATION
# =============================================================================

def create_title_slide(prs: Presentation, slide_data: dict, palette: dict):
    """Create a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_pptx_rgb(palette["primary"])
    background.line.fill.background()

    # Title
    title = slide_data.get("title", "Untitled Presentation")
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2),
        Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_pptx_rgb(palette["background"])
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.7),
            Inches(9), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = hex_to_pptx_rgb(palette["secondary"])
        subtitle_para.alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide(prs: Presentation, slide_data: dict, palette: dict):
    """Create a content slide with title and body/bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = hex_to_pptx_rgb(palette["primary"])
    header_bar.line.fill.background()

    # Title
    title = slide_data.get("title", "")
    if title:
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = hex_to_pptx_rgb(palette["background"])

    # Body content or bullets
    content_top = Inches(1.5)
    bullets = slide_data.get("bullets", [])
    body = slide_data.get("body", "")

    if bullets:
        content_box = slide.shapes.add_textbox(
            Inches(0.75), content_top,
            Inches(8.5), Inches(5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()
            para.text = f"  {bullet}"
            para.font.size = Pt(20)
            para.font.color.rgb = hex_to_pptx_rgb(palette["text"])
            para.space_before = Pt(12)
            para.level = 0

    elif body:
        content_box = slide.shapes.add_textbox(
            Inches(0.75), content_top,
            Inches(8.5), Inches(5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.paragraphs[0].text = body
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.color.rgb = hex_to_pptx_rgb(palette["text"])

    return slide

def create_two_column_slide(prs: Presentation, slide_data: dict, palette: dict):
    """Create a two-column slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = hex_to_pptx_rgb(palette["primary"])
    header_bar.line.fill.background()

    # Title
    title = slide_data.get("title", "")
    if title:
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = hex_to_pptx_rgb(palette["background"])

    # Left column
    left_content = slide_data.get("left_content", "")
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(4.25), Inches(5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_frame.paragraphs[0].text = left_content
    left_frame.paragraphs[0].font.size = Pt(16)
    left_frame.paragraphs[0].font.color.rgb = hex_to_pptx_rgb(palette["text"])

    # Right column
    right_content = slide_data.get("right_content", "")
    right_box = slide.shapes.add_textbox(
        Inches(5.25), Inches(1.5),
        Inches(4.25), Inches(5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_frame.paragraphs[0].text = right_content
    right_frame.paragraphs[0].font.size = Pt(16)
    right_frame.paragraphs[0].font.color.rgb = hex_to_pptx_rgb(palette["text"])

    # Divider line
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.875), Inches(1.5),
        Inches(0.02), Inches(5)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_pptx_rgb(palette["accent"])
    divider.line.fill.background()

    return slide

def create_section_slide(prs: Presentation, slide_data: dict, palette: dict):
    """Create a section header slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Left accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.3), prs.slide_height
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_pptx_rgb(palette["secondary"])
    accent_bar.line.fill.background()

    # Title
    title = slide_data.get("title", "Section")
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.5),
        Inches(8.5), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_pptx_rgb(palette["primary"])

    # Subtitle
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(4.2),
            Inches(8.5), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = hex_to_pptx_rgb(palette["accent"])

    return slide

def create_quote_slide(prs: Presentation, slide_data: dict, palette: dict):
    """Create a quote slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background with subtle color
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_pptx_rgb(palette["light"])
    background.line.fill.background()

    # Quote mark decoration
    quote_mark = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(1), Inches(1)
    )
    quote_frame = quote_mark.text_frame
    quote_para = quote_frame.paragraphs[0]
    quote_para.text = '"'
    quote_para.font.size = Pt(120)
    quote_para.font.color.rgb = hex_to_pptx_rgb(palette["secondary"])

    # Quote text
    quote = slide_data.get("body", slide_data.get("title", ""))
    quote_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(3)
    )
    quote_text_frame = quote_box.text_frame
    quote_text_frame.word_wrap = True
    quote_text_para = quote_text_frame.paragraphs[0]
    quote_text_para.text = quote
    quote_text_para.font.size = Pt(28)
    quote_text_para.font.italic = True
    quote_text_para.font.color.rgb = hex_to_pptx_rgb(palette["primary"])
    quote_text_para.alignment = PP_ALIGN.CENTER

    # Attribution
    attribution = slide_data.get("subtitle", "")
    if attribution:
        attr_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.5),
            Inches(8), Inches(0.5)
        )
        attr_frame = attr_box.text_frame
        attr_para = attr_frame.paragraphs[0]
        attr_para.text = f"- {attribution}"
        attr_para.font.size = Pt(18)
        attr_para.font.color.rgb = hex_to_pptx_rgb(palette["accent"])
        attr_para.alignment = PP_ALIGN.RIGHT

    return slide

def generate_presentation(content: Dict[str, Any], palette: dict) -> str:
    """Generate a PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slides_data = content.get("slides", [])

    for slide_data in slides_data:
        layout = slide_data.get("layout", "content")

        if layout == "title":
            create_title_slide(prs, slide_data, palette)
        elif layout == "content" or layout == "bullet_points":
            create_content_slide(prs, slide_data, palette)
        elif layout == "two_column":
            create_two_column_slide(prs, slide_data, palette)
        elif layout == "section":
            create_section_slide(prs, slide_data, palette)
        elif layout == "quote":
            create_quote_slide(prs, slide_data, palette)
        else:
            # Default to content slide
            create_content_slide(prs, slide_data, palette)

    # Save presentation
    file_id = generate_file_id()
    output_path = f"{OUTPUT_DIR}/presentation_{file_id}.pptx"
    prs.save(output_path)

    return output_path

# =============================================================================
# DOCUMENT GENERATION
# =============================================================================

def generate_document(content: str, format: str, palette: dict,
                     title: str = None, author: str = None,
                     headers: dict = None) -> str:
    """Generate a Word document with brand styling"""
    doc = Document()

    # Set document properties
    core_props = doc.core_properties
    if title:
        core_props.title = title
    if author:
        core_props.author = author

    # Add title if provided
    if title:
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in title_para.runs:
            run.font.color.rgb = hex_to_docx_rgb(palette["primary"])

    # Process content - split by headers (## or lines starting with headers)
    lines = content.split('\n')
    current_para = None

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # Check for markdown-style headers
        if line.startswith('# '):
            heading = doc.add_heading(line[2:], level=1)
            for run in heading.runs:
                run.font.color.rgb = hex_to_docx_rgb(palette["primary"])
        elif line.startswith('## '):
            heading = doc.add_heading(line[3:], level=2)
            for run in heading.runs:
                run.font.color.rgb = hex_to_docx_rgb(palette["primary"])
        elif line.startswith('### '):
            heading = doc.add_heading(line[4:], level=3)
            for run in heading.runs:
                run.font.color.rgb = hex_to_docx_rgb(palette["accent"])
        elif line.startswith('- ') or line.startswith('* '):
            # Bullet point
            para = doc.add_paragraph(line[2:], style='List Bullet')
            for run in para.runs:
                run.font.color.rgb = hex_to_docx_rgb(palette["text"])
        elif line.startswith('1. ') or (len(line) > 2 and line[0].isdigit() and line[1] == '.'):
            # Numbered list
            para = doc.add_paragraph(line.split('. ', 1)[1] if '. ' in line else line,
                                    style='List Number')
            for run in para.runs:
                run.font.color.rgb = hex_to_docx_rgb(palette["text"])
        else:
            # Regular paragraph
            para = doc.add_paragraph(line)
            for run in para.runs:
                run.font.color.rgb = hex_to_docx_rgb(palette["text"])

    # Save document
    file_id = generate_file_id()

    if format == "docx":
        output_path = f"{OUTPUT_DIR}/document_{file_id}.docx"
        doc.save(output_path)
    else:
        # For PDF, we'll save as docx (PDF conversion would require additional libraries)
        # In production, use python-docx2pdf or similar
        output_path = f"{OUTPUT_DIR}/document_{file_id}.docx"
        doc.save(output_path)

    return output_path

# =============================================================================
# CHART GENERATION
# =============================================================================

def generate_chart(chart_type: str, data: Dict[str, Any], style: str,
                   title: str = None, width: int = 800, height: int = 600) -> dict:
    """Generate a chart as PNG and SVG"""
    palette = get_palette(style)

    # Set up matplotlib style
    plt.style.use('seaborn-v0_8-whitegrid')
    fig, ax = plt.subplots(figsize=(width/100, height/100), dpi=100)

    labels = data.get("labels", [])
    values = data.get("values", [])
    series = data.get("series", None)

    # Color cycle from palette
    colors = [
        palette["primary"],
        palette["secondary"],
        palette["accent"],
        palette["success"],
        palette["warning"],
        palette["error"]
    ]

    if chart_type == "bar":
        if series:
            # Multi-series bar chart
            x = np.arange(len(labels))
            width_bar = 0.8 / len(series)
            for i, s in enumerate(series):
                offset = (i - len(series)/2 + 0.5) * width_bar
                ax.bar(x + offset, s["values"], width_bar, label=s.get("name", f"Series {i+1}"),
                      color=colors[i % len(colors)])
            ax.set_xticks(x)
            ax.set_xticklabels(labels)
            ax.legend()
        else:
            ax.bar(labels, values, color=colors[0])

    elif chart_type == "hbar":
        ax.barh(labels, values, color=colors[0])

    elif chart_type == "line":
        if series:
            for i, s in enumerate(series):
                ax.plot(labels, s["values"], marker='o', label=s.get("name", f"Series {i+1}"),
                       color=colors[i % len(colors)], linewidth=2)
            ax.legend()
        else:
            ax.plot(labels, values, marker='o', color=colors[0], linewidth=2)

    elif chart_type == "area":
        if series:
            for i, s in enumerate(series):
                ax.fill_between(range(len(labels)), s["values"], alpha=0.6,
                               label=s.get("name", f"Series {i+1}"), color=colors[i % len(colors)])
            ax.set_xticks(range(len(labels)))
            ax.set_xticklabels(labels)
            ax.legend()
        else:
            ax.fill_between(range(len(labels)), values, alpha=0.6, color=colors[0])
            ax.set_xticks(range(len(labels)))
            ax.set_xticklabels(labels)

    elif chart_type == "pie":
        ax.pie(values, labels=labels, autopct='%1.1f%%', colors=colors[:len(values)],
               startangle=90)
        ax.axis('equal')

    elif chart_type == "donut":
        wedges, texts, autotexts = ax.pie(values, labels=labels, autopct='%1.1f%%',
                                          colors=colors[:len(values)], startangle=90,
                                          pctdistance=0.75)
        # Draw center circle for donut effect
        centre_circle = plt.Circle((0, 0), 0.50, fc='white')
        ax.add_patch(centre_circle)
        ax.axis('equal')

    elif chart_type == "scatter":
        x_values = data.get("x", values)
        y_values = data.get("y", values)
        ax.scatter(x_values, y_values, c=colors[0], s=100, alpha=0.7)

    # Title
    if title:
        ax.set_title(title, fontsize=16, fontweight='bold', color=palette["text"])

    # Style adjustments
    ax.tick_params(colors=palette["text"])
    for spine in ax.spines.values():
        spine.set_color(palette["accent"])

    plt.tight_layout()

    # Save as PNG and SVG
    file_id = generate_file_id()
    png_path = f"{OUTPUT_DIR}/chart_{file_id}.png"
    svg_path = f"{OUTPUT_DIR}/chart_{file_id}.svg"

    fig.savefig(png_path, dpi=150, bbox_inches='tight', facecolor='white')
    fig.savefig(svg_path, format='svg', bbox_inches='tight', facecolor='white')

    plt.close(fig)

    return {
        "image_path": png_path,
        "svg_path": svg_path,
        "animation_path": None  # Animation would require additional libraries
    }

# =============================================================================
# THUMBNAIL GENERATION
# =============================================================================

def create_gradient_background(size: tuple, color1: str, color2: str, direction: str = "diagonal") -> Image.Image:
    """Create a gradient background"""
    width, height = size
    img = Image.new('RGB', size)
    draw = ImageDraw.Draw(img)

    r1, g1, b1 = hex_to_rgb(color1)
    r2, g2, b2 = hex_to_rgb(color2)

    if direction == "diagonal":
        for y in range(height):
            for x in range(width):
                # Diagonal gradient
                ratio = (x + y) / (width + height)
                r = int(r1 + (r2 - r1) * ratio)
                g = int(g1 + (g2 - g1) * ratio)
                b = int(b1 + (b2 - b1) * ratio)
                draw.point((x, y), fill=(r, g, b))
    elif direction == "horizontal":
        for x in range(width):
            ratio = x / width
            r = int(r1 + (r2 - r1) * ratio)
            g = int(g1 + (g2 - g1) * ratio)
            b = int(b1 + (b2 - b1) * ratio)
            draw.line([(x, 0), (x, height)], fill=(r, g, b))
    else:  # vertical
        for y in range(height):
            ratio = y / height
            r = int(r1 + (r2 - r1) * ratio)
            g = int(g1 + (g2 - g1) * ratio)
            b = int(b1 + (b2 - b1) * ratio)
            draw.line([(0, y), (width, y)], fill=(r, g, b))

    return img

def generate_thumbnail(title: str, style: str, palette: dict,
                       background_image: str = None,
                       subtitle: str = None) -> dict:
    """Generate thumbnails in multiple sizes"""
    sizes = {
        "small": (320, 180),
        "medium": (640, 360),
        "large": (1280, 720)
    }

    file_id = generate_file_id()
    output_paths = {}

    for size_name, dimensions in sizes.items():
        width, height = dimensions

        # Create background
        if style == "gradient":
            img = create_gradient_background(dimensions, palette["primary"], palette["accent"])
        elif style == "solid":
            img = Image.new('RGB', dimensions, hex_to_rgb(palette["primary"]))
        elif style == "pattern":
            # Simple pattern with alternating colors
            img = Image.new('RGB', dimensions, hex_to_rgb(palette["background"]))
            draw = ImageDraw.Draw(img)
            pattern_color = hex_to_rgb(palette["light"])
            for i in range(0, width, 20):
                for j in range(0, height, 20):
                    if (i + j) % 40 == 0:
                        draw.rectangle([i, j, i+10, j+10], fill=pattern_color)
        elif background_image and style == "image":
            try:
                img = Image.open(background_image)
                img = img.resize(dimensions, Image.Resampling.LANCZOS)
                # Add overlay for text readability
                overlay = Image.new('RGBA', dimensions, (*hex_to_rgb(palette["primary"]), 180))
                img = img.convert('RGBA')
                img = Image.alpha_composite(img, overlay)
                img = img.convert('RGB')
            except Exception:
                img = create_gradient_background(dimensions, palette["primary"], palette["accent"])
        else:
            img = create_gradient_background(dimensions, palette["primary"], palette["accent"])

        draw = ImageDraw.Draw(img)

        # Try to load a good font, fall back to default
        try:
            # Calculate font size based on image dimensions
            title_font_size = int(height * 0.12)
            subtitle_font_size = int(height * 0.06)
            title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", title_font_size)
            subtitle_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", subtitle_font_size)
        except Exception:
            title_font = ImageFont.load_default()
            subtitle_font = ImageFont.load_default()

        # Draw title (centered)
        text_color = hex_to_rgb(palette["background"])

        # Word wrap title if too long
        max_chars = int(width / (title_font_size * 0.6))
        if len(title) > max_chars:
            words = title.split()
            lines = []
            current_line = ""
            for word in words:
                if len(current_line) + len(word) + 1 <= max_chars:
                    current_line += (" " if current_line else "") + word
                else:
                    if current_line:
                        lines.append(current_line)
                    current_line = word
            if current_line:
                lines.append(current_line)
            title_text = "\n".join(lines)
        else:
            title_text = title

        # Calculate text position
        bbox = draw.textbbox((0, 0), title_text, font=title_font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        title_y = (height - text_height) // 2 - (int(height * 0.05) if subtitle else 0)
        title_x = (width - text_width) // 2

        draw.text((title_x, title_y), title_text, font=title_font, fill=text_color, align="center")

        # Draw subtitle if provided
        if subtitle:
            bbox = draw.textbbox((0, 0), subtitle, font=subtitle_font)
            sub_width = bbox[2] - bbox[0]
            sub_x = (width - sub_width) // 2
            sub_y = title_y + text_height + int(height * 0.03)
            sub_color = hex_to_rgb(palette["secondary"])
            draw.text((sub_x, sub_y), subtitle, font=subtitle_font, fill=sub_color)

        # Save
        output_path = f"{OUTPUT_DIR}/thumbnail_{file_id}_{size_name}.png"
        img.save(output_path, "PNG")
        output_paths[size_name] = output_path

    return output_paths

# =============================================================================
# WEB/UI DESIGN GENERATION
# =============================================================================

# Tailwind-based style configurations
WEB_STYLES = {
    "modern": {
        "font_class": "font-sans",
        "heading_class": "font-bold tracking-tight",
        "container_class": "max-w-7xl mx-auto px-4 sm:px-6 lg:px-8",
        "section_padding": "py-16 sm:py-24",
        "button_class": "rounded-lg shadow-lg hover:shadow-xl transition-all duration-300",
        "card_class": "rounded-2xl shadow-lg hover:shadow-xl transition-shadow",
        "gradient_direction": "bg-gradient-to-br",
    },
    "minimal": {
        "font_class": "font-sans",
        "heading_class": "font-medium tracking-normal",
        "container_class": "max-w-5xl mx-auto px-4",
        "section_padding": "py-12 sm:py-16",
        "button_class": "rounded-none border-2 hover:bg-black hover:text-white transition-colors",
        "card_class": "border border-gray-200",
        "gradient_direction": "bg-gradient-to-r",
    },
    "corporate": {
        "font_class": "font-serif",
        "heading_class": "font-bold tracking-normal",
        "container_class": "max-w-6xl mx-auto px-6",
        "section_padding": "py-16 sm:py-20",
        "button_class": "rounded-md shadow hover:shadow-md transition-shadow",
        "card_class": "rounded-lg shadow-md border border-gray-100",
        "gradient_direction": "bg-gradient-to-b",
    },
    "playful": {
        "font_class": "font-sans",
        "heading_class": "font-extrabold tracking-tight",
        "container_class": "max-w-7xl mx-auto px-4 sm:px-8",
        "section_padding": "py-20 sm:py-28",
        "button_class": "rounded-full shadow-lg hover:scale-105 transition-transform",
        "card_class": "rounded-3xl shadow-xl hover:rotate-1 transition-transform",
        "gradient_direction": "bg-gradient-to-tr",
    },
    "dark": {
        "font_class": "font-sans",
        "heading_class": "font-bold tracking-tight",
        "container_class": "max-w-7xl mx-auto px-4 sm:px-6 lg:px-8",
        "section_padding": "py-16 sm:py-24",
        "button_class": "rounded-lg shadow-lg hover:shadow-2xl transition-all duration-300 border border-gray-700",
        "card_class": "rounded-2xl bg-gray-800 border border-gray-700",
        "gradient_direction": "bg-gradient-to-br",
    }
}

def get_tailwind_colors(palette: dict, style: str = "modern") -> dict:
    """Convert palette to Tailwind-compatible color classes"""
    # Map hex colors to closest Tailwind approximations for inline styles
    return {
        "primary": palette["primary"],
        "secondary": palette["secondary"],
        "accent": palette["accent"],
        "background": palette["background"],
        "text": palette["text"],
        "light": palette.get("light", "#F5F7FA"),
    }

def generate_tailwind_base() -> str:
    """Generate base HTML with Tailwind CDN"""
    return '''<!DOCTYPE html>
<html lang="en" class="scroll-smooth">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <script src="https://cdn.tailwindcss.com"></script>
    <link rel="preconnect" href="https://fonts.googleapis.com">
    <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap" rel="stylesheet">
    {favicon}
    <script>
        tailwind.config = {{
            theme: {{
                extend: {{
                    colors: {{
                        primary: '{primary}',
                        secondary: '{secondary}',
                        accent: '{accent}',
                    }},
                    fontFamily: {{
                        sans: ['Inter', 'system-ui', 'sans-serif'],
                    }}
                }}
            }}
        }}
    </script>
    <style>
        .animate-fade-in {{ animation: fadeIn 0.6s ease-out; }}
        .animate-slide-up {{ animation: slideUp 0.6s ease-out; }}
        @keyframes fadeIn {{ from {{ opacity: 0; }} to {{ opacity: 1; }} }}
        @keyframes slideUp {{ from {{ opacity: 0; transform: translateY(20px); }} to {{ opacity: 1; transform: translateY(0); }} }}
    </style>
</head>
<body class="{body_class}">
'''

def generate_hero_section(section: ContentSection, palette: dict, style_config: dict, animate: bool = True) -> str:
    """Generate a hero section"""
    animation = "animate-fade-in" if animate else ""
    gradient = f"from-[{palette['primary']}] to-[{palette['accent']}]" if section.background == "gradient" else ""

    bg_style = ""
    if section.background and section.background.startswith("#"):
        bg_style = f'style="background-color: {section.background};"'
    elif section.background == "gradient":
        bg_style = f'style="background: linear-gradient(135deg, {palette["primary"]}, {palette["accent"]});"'

    return f'''
    <section class="relative {style_config['section_padding']} overflow-hidden {animation}" {bg_style}>
        <div class="{style_config['container_class']}">
            <div class="text-center max-w-4xl mx-auto">
                <h1 class="text-4xl sm:text-5xl lg:text-6xl {style_config['heading_class']} mb-6" style="color: {palette['text'] if not section.background else palette['background']};">
                    {section.title or "Welcome"}
                </h1>
                <p class="text-xl sm:text-2xl mb-8 opacity-80" style="color: {palette['text'] if not section.background else palette['background']};">
                    {section.subtitle or ""}
                </p>
                {f'<p class="text-lg mb-10 opacity-70" style="color: {palette["text"] if not section.background else palette["background"]};">{section.content}</p>' if section.content else ''}
                {f'''<a href="{section.cta_url}" class="inline-flex items-center px-8 py-4 text-lg font-semibold {style_config['button_class']}" style="background-color: {palette['secondary']}; color: {palette['background']};">
                    {section.cta_text}
                    <svg class="ml-2 w-5 h-5" fill="none" stroke="currentColor" viewBox="0 0 24 24"><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M17 8l4 4m0 0l-4 4m4-4H3"></path></svg>
                </a>''' if section.cta_text else ''}
            </div>
            {f'<img src="{section.image_url}" alt="" class="mt-12 mx-auto rounded-xl shadow-2xl max-w-full">' if section.image_url else ''}
        </div>
    </section>
'''

def generate_features_section(section: ContentSection, palette: dict, style_config: dict, animate: bool = True) -> str:
    """Generate a features section"""
    animation = "animate-slide-up" if animate else ""
    features = section.items or []

    features_html = ""
    for i, feature in enumerate(features):
        delay = f"animation-delay: {i * 0.1}s;" if animate else ""
        features_html += f'''
            <div class="{style_config['card_class']} p-6 {animation}" style="{delay}">
                {f'<div class="w-12 h-12 mb-4 flex items-center justify-center rounded-lg" style="background-color: {palette["secondary"]}20;"><span class="text-2xl">{feature.get("icon", "✨")}</span></div>' if feature.get("icon") else ''}
                <h3 class="text-xl font-semibold mb-2" style="color: {palette['text']};">{feature.get("title", "Feature")}</h3>
                <p class="opacity-70" style="color: {palette['text']};">{feature.get("description", "")}</p>
            </div>
        '''

    return f'''
    <section class="{style_config['section_padding']}" style="background-color: {palette['light']};">
        <div class="{style_config['container_class']}">
            <div class="text-center mb-12">
                <h2 class="text-3xl sm:text-4xl {style_config['heading_class']} mb-4" style="color: {palette['text']};">{section.title or "Features"}</h2>
                <p class="text-xl opacity-70 max-w-2xl mx-auto" style="color: {palette['text']};">{section.subtitle or ""}</p>
            </div>
            <div class="grid grid-cols-1 md:grid-cols-2 lg:grid-cols-3 gap-8">
                {features_html}
            </div>
        </div>
    </section>
'''

def generate_cta_section(section: ContentSection, palette: dict, style_config: dict, animate: bool = True) -> str:
    """Generate a CTA section"""
    animation = "animate-fade-in" if animate else ""

    return f'''
    <section class="{style_config['section_padding']} {animation}" style="background: linear-gradient(135deg, {palette['primary']}, {palette['accent']});">
        <div class="{style_config['container_class']} text-center">
            <h2 class="text-3xl sm:text-4xl {style_config['heading_class']} mb-4" style="color: {palette['background']};">{section.title or "Ready to get started?"}</h2>
            <p class="text-xl mb-8 opacity-90" style="color: {palette['background']};">{section.subtitle or ""}</p>
            <a href="{section.cta_url}" class="inline-flex items-center px-8 py-4 text-lg font-semibold {style_config['button_class']}" style="background-color: {palette['background']}; color: {palette['primary']};">
                {section.cta_text or "Get Started"}
            </a>
        </div>
    </section>
'''

def generate_testimonials_section(section: ContentSection, palette: dict, style_config: dict, animate: bool = True) -> str:
    """Generate a testimonials section"""
    animation = "animate-slide-up" if animate else ""
    testimonials = section.items or []

    testimonials_html = ""
    for i, testimonial in enumerate(testimonials):
        delay = f"animation-delay: {i * 0.15}s;" if animate else ""
        testimonials_html += f'''
            <div class="{style_config['card_class']} p-8 {animation}" style="background-color: {palette['background']}; {delay}">
                <div class="flex items-center mb-4">
                    <div class="w-12 h-12 rounded-full flex items-center justify-center text-xl font-bold" style="background-color: {palette['primary']}; color: {palette['background']};">
                        {testimonial.get("name", "A")[0]}
                    </div>
                    <div class="ml-4">
                        <p class="font-semibold" style="color: {palette['text']};">{testimonial.get("name", "Customer")}</p>
                        <p class="text-sm opacity-60" style="color: {palette['text']};">{testimonial.get("role", "")}</p>
                    </div>
                </div>
                <p class="italic opacity-80" style="color: {palette['text']};">"{testimonial.get("quote", "")}"</p>
            </div>
        '''

    return f'''
    <section class="{style_config['section_padding']}" style="background-color: {palette['light']};">
        <div class="{style_config['container_class']}">
            <div class="text-center mb-12">
                <h2 class="text-3xl sm:text-4xl {style_config['heading_class']} mb-4" style="color: {palette['text']};">{section.title or "What People Say"}</h2>
                <p class="text-xl opacity-70" style="color: {palette['text']};">{section.subtitle or ""}</p>
            </div>
            <div class="grid grid-cols-1 md:grid-cols-2 lg:grid-cols-3 gap-8">
                {testimonials_html}
            </div>
        </div>
    </section>
'''

def generate_pricing_section(section: ContentSection, palette: dict, style_config: dict, animate: bool = True) -> str:
    """Generate a pricing section"""
    animation = "animate-slide-up" if animate else ""
    plans = section.items or []

    plans_html = ""
    for i, plan in enumerate(plans):
        is_featured = plan.get("featured", False)
        delay = f"animation-delay: {i * 0.1}s;" if animate else ""
        border = f"border-2" if is_featured else ""
        border_color = f"border-[{palette['secondary']}]" if is_featured else "border-gray-200"
        scale = "scale-105" if is_featured else ""

        features_list = ""
        for feature in plan.get("features", []):
            features_list += f'<li class="flex items-center py-2"><svg class="w-5 h-5 mr-3" style="color: {palette["secondary"]};" fill="currentColor" viewBox="0 0 20 20"><path fill-rule="evenodd" d="M16.707 5.293a1 1 0 010 1.414l-8 8a1 1 0 01-1.414 0l-4-4a1 1 0 011.414-1.414L8 12.586l7.293-7.293a1 1 0 011.414 0z" clip-rule="evenodd"></path></svg>{feature}</li>'

        plans_html += f'''
            <div class="{style_config['card_class']} p-8 {border} {scale} {animation}" style="background-color: {palette['background']}; {f'border-color: {palette["secondary"]};' if is_featured else ''} {delay}">
                {f'<span class="inline-block px-3 py-1 text-sm font-semibold rounded-full mb-4" style="background-color: {palette["secondary"]}; color: {palette["background"]};">Most Popular</span>' if is_featured else ''}
                <h3 class="text-2xl font-bold mb-2" style="color: {palette['text']};">{plan.get("name", "Plan")}</h3>
                <p class="opacity-60 mb-4" style="color: {palette['text']};">{plan.get("description", "")}</p>
                <div class="mb-6">
                    <span class="text-4xl font-bold" style="color: {palette['text']};">{plan.get("price", "$0")}</span>
                    <span class="opacity-60" style="color: {palette['text']};">/{plan.get("period", "month")}</span>
                </div>
                <ul class="mb-8 space-y-1" style="color: {palette['text']};">
                    {features_list}
                </ul>
                <a href="{plan.get('cta_url', '#')}" class="block w-full text-center py-3 {style_config['button_class']}" style="background-color: {palette['secondary'] if is_featured else palette['primary']}; color: {palette['background']};">
                    {plan.get("cta_text", "Get Started")}
                </a>
            </div>
        '''

    return f'''
    <section class="{style_config['section_padding']}" style="background-color: {palette['background']};">
        <div class="{style_config['container_class']}">
            <div class="text-center mb-12">
                <h2 class="text-3xl sm:text-4xl {style_config['heading_class']} mb-4" style="color: {palette['text']};">{section.title or "Pricing"}</h2>
                <p class="text-xl opacity-70 max-w-2xl mx-auto" style="color: {palette['text']};">{section.subtitle or ""}</p>
            </div>
            <div class="grid grid-cols-1 md:grid-cols-2 lg:grid-cols-3 gap-8 items-center">
                {plans_html}
            </div>
        </div>
    </section>
'''

def generate_stats_section(section: ContentSection, palette: dict, style_config: dict, animate: bool = True) -> str:
    """Generate a stats section"""
    animation = "animate-fade-in" if animate else ""
    stats = section.items or []

    stats_html = ""
    for stat in stats:
        stats_html += f'''
            <div class="text-center">
                <p class="text-4xl sm:text-5xl font-bold mb-2" style="color: {palette['secondary']};">{stat.get("value", "0")}</p>
                <p class="text-lg opacity-70" style="color: {palette['text']};">{stat.get("label", "")}</p>
            </div>
        '''

    return f'''
    <section class="{style_config['section_padding']} {animation}" style="background-color: {palette['primary']};">
        <div class="{style_config['container_class']}">
            <div class="text-center mb-12">
                <h2 class="text-3xl sm:text-4xl {style_config['heading_class']} mb-4" style="color: {palette['background']};">{section.title or ""}</h2>
            </div>
            <div class="grid grid-cols-2 md:grid-cols-4 gap-8">
                {stats_html}
            </div>
        </div>
    </section>
'''

def generate_footer_section(section: ContentSection, palette: dict, style_config: dict, animate: bool = True) -> str:
    """Generate a footer section"""
    links = section.items or []

    links_html = ""
    for group in links:
        group_links = ""
        for link in group.get("links", []):
            group_links += f'<li><a href="{link.get("url", "#")}" class="hover:opacity-100 opacity-70 transition-opacity">{link.get("text", "")}</a></li>'
        links_html += f'''
            <div>
                <h4 class="font-semibold mb-4">{group.get("title", "")}</h4>
                <ul class="space-y-2">
                    {group_links}
                </ul>
            </div>
        '''

    return f'''
    <footer class="py-12" style="background-color: {palette['primary']}; color: {palette['background']};">
        <div class="{style_config['container_class']}">
            <div class="grid grid-cols-2 md:grid-cols-4 gap-8 mb-8">
                {links_html}
            </div>
            <div class="pt-8 border-t border-white/20 text-center opacity-70">
                <p>{section.content or f"&copy; {datetime.now().year} All rights reserved."}</p>
            </div>
        </div>
    </footer>
'''

def generate_faq_section(section: ContentSection, palette: dict, style_config: dict, animate: bool = True) -> str:
    """Generate an FAQ section"""
    animation = "animate-slide-up" if animate else ""
    faqs = section.items or []

    faqs_html = ""
    for i, faq in enumerate(faqs):
        faqs_html += f'''
            <div class="{style_config['card_class']} p-6 mb-4" style="background-color: {palette['background']};">
                <details class="group">
                    <summary class="flex justify-between items-center cursor-pointer list-none">
                        <h3 class="font-semibold text-lg" style="color: {palette['text']};">{faq.get("question", "")}</h3>
                        <span class="transition-transform group-open:rotate-180">
                            <svg class="w-5 h-5" fill="none" stroke="currentColor" viewBox="0 0 24 24"><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M19 9l-7 7-7-7"></path></svg>
                        </span>
                    </summary>
                    <p class="mt-4 opacity-70" style="color: {palette['text']};">{faq.get("answer", "")}</p>
                </details>
            </div>
        '''

    return f'''
    <section class="{style_config['section_padding']} {animation}" style="background-color: {palette['light']};">
        <div class="{style_config['container_class']}">
            <div class="text-center mb-12">
                <h2 class="text-3xl sm:text-4xl {style_config['heading_class']} mb-4" style="color: {palette['text']};">{section.title or "Frequently Asked Questions"}</h2>
                <p class="text-xl opacity-70" style="color: {palette['text']};">{section.subtitle or ""}</p>
            </div>
            <div class="max-w-3xl mx-auto">
                {faqs_html}
            </div>
        </div>
    </section>
'''

def generate_contact_section(section: ContentSection, palette: dict, style_config: dict, animate: bool = True) -> str:
    """Generate a contact section"""
    animation = "animate-fade-in" if animate else ""

    return f'''
    <section class="{style_config['section_padding']} {animation}" style="background-color: {palette['background']};">
        <div class="{style_config['container_class']}">
            <div class="grid grid-cols-1 lg:grid-cols-2 gap-12">
                <div>
                    <h2 class="text-3xl sm:text-4xl {style_config['heading_class']} mb-4" style="color: {palette['text']};">{section.title or "Get in Touch"}</h2>
                    <p class="text-xl opacity-70 mb-8" style="color: {palette['text']};">{section.subtitle or ""}</p>
                    <div class="space-y-4">
                        <p class="flex items-center" style="color: {palette['text']};"><svg class="w-6 h-6 mr-3" style="color: {palette['secondary']};" fill="none" stroke="currentColor" viewBox="0 0 24 24"><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M3 8l7.89 5.26a2 2 0 002.22 0L21 8M5 19h14a2 2 0 002-2V7a2 2 0 00-2-2H5a2 2 0 00-2 2v10a2 2 0 002 2z"></path></svg>hello@example.com</p>
                        <p class="flex items-center" style="color: {palette['text']};"><svg class="w-6 h-6 mr-3" style="color: {palette['secondary']};" fill="none" stroke="currentColor" viewBox="0 0 24 24"><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M17.657 16.657L13.414 20.9a1.998 1.998 0 01-2.827 0l-4.244-4.243a8 8 0 1111.314 0z"></path><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M15 11a3 3 0 11-6 0 3 3 0 016 0z"></path></svg>123 Main St, City</p>
                    </div>
                </div>
                <div class="{style_config['card_class']} p-8" style="background-color: {palette['light']};">
                    <form class="space-y-6">
                        <div>
                            <label class="block text-sm font-medium mb-2" style="color: {palette['text']};">Name</label>
                            <input type="text" class="w-full px-4 py-3 rounded-lg border focus:outline-none focus:ring-2" style="border-color: {palette['accent']}33; --tw-ring-color: {palette['secondary']};">
                        </div>
                        <div>
                            <label class="block text-sm font-medium mb-2" style="color: {palette['text']};">Email</label>
                            <input type="email" class="w-full px-4 py-3 rounded-lg border focus:outline-none focus:ring-2" style="border-color: {palette['accent']}33; --tw-ring-color: {palette['secondary']};">
                        </div>
                        <div>
                            <label class="block text-sm font-medium mb-2" style="color: {palette['text']};">Message</label>
                            <textarea rows="4" class="w-full px-4 py-3 rounded-lg border focus:outline-none focus:ring-2" style="border-color: {palette['accent']}33; --tw-ring-color: {palette['secondary']};"></textarea>
                        </div>
                        <button type="submit" class="w-full py-3 {style_config['button_class']}" style="background-color: {palette['primary']}; color: {palette['background']};">{section.cta_text or "Send Message"}</button>
                    </form>
                </div>
            </div>
        </div>
    </section>
'''

def generate_landing_page(req: LandingPageRequest, palette: dict) -> str:
    """Generate a complete landing page"""
    style_config = WEB_STYLES.get(req.style, WEB_STYLES["modern"])

    # Apply custom colors if provided
    if req.custom_colors:
        palette = {**palette, **req.custom_colors}

    # Generate base HTML
    body_class = "bg-white text-gray-900" if req.style != "dark" else "bg-gray-900 text-white"
    favicon = f'<link rel="icon" href="{req.favicon_url}">' if req.favicon_url else ''

    html = generate_tailwind_base().format(
        title=req.title,
        primary=palette["primary"],
        secondary=palette["secondary"],
        accent=palette["accent"],
        body_class=body_class,
        favicon=favicon
    )

    # Generate navbar if logo provided
    if req.logo_url:
        html += f'''
    <nav class="fixed top-0 left-0 right-0 z-50 py-4" style="background-color: {palette['background']}ee; backdrop-filter: blur(10px);">
        <div class="{style_config['container_class']} flex justify-between items-center">
            <img src="{req.logo_url}" alt="Logo" class="h-8">
            <div class="hidden md:flex space-x-8">
                <a href="#" class="hover:opacity-70 transition-opacity" style="color: {palette['text']};">Home</a>
                <a href="#features" class="hover:opacity-70 transition-opacity" style="color: {palette['text']};">Features</a>
                <a href="#pricing" class="hover:opacity-70 transition-opacity" style="color: {palette['text']};">Pricing</a>
                <a href="#contact" class="hover:opacity-70 transition-opacity" style="color: {palette['text']};">Contact</a>
            </div>
            <button class="px-6 py-2 {style_config['button_class']}" style="background-color: {palette['primary']}; color: {palette['background']};">Get Started</button>
        </div>
    </nav>
    <div class="h-20"></div>
'''

    # Section generators
    section_generators = {
        "hero": generate_hero_section,
        "features": generate_features_section,
        "cta": generate_cta_section,
        "testimonials": generate_testimonials_section,
        "pricing": generate_pricing_section,
        "stats": generate_stats_section,
        "footer": generate_footer_section,
        "faq": generate_faq_section,
        "contact": generate_contact_section,
    }

    # Generate each section
    for section in req.sections:
        generator = section_generators.get(section.type)
        if generator:
            html += generator(section, palette, style_config, req.include_animations)
        elif section.type == "custom" and section.content:
            html += f'''
    <section class="{style_config['section_padding']}">
        <div class="{style_config['container_class']}">
            {section.content}
        </div>
    </section>
'''

    html += '''
</body>
</html>'''

    return html

def generate_ui_component(req: UIComponentRequest, palette: dict) -> str:
    """Generate a UI component"""
    style_config = WEB_STYLES.get("modern", WEB_STYLES["modern"])
    content = req.content or {}

    # Size mappings
    sizes = {
        "sm": {"padding": "px-3 py-1.5", "text": "text-sm", "icon": "w-4 h-4"},
        "md": {"padding": "px-4 py-2", "text": "text-base", "icon": "w-5 h-5"},
        "lg": {"padding": "px-6 py-3", "text": "text-lg", "icon": "w-6 h-6"},
        "xl": {"padding": "px-8 py-4", "text": "text-xl", "icon": "w-7 h-7"},
    }
    size = sizes.get(req.size, sizes["md"])

    # Variant colors
    if req.variant == "primary":
        bg = palette["primary"]
        text_color = palette["background"]
        border = palette["primary"]
    elif req.variant == "secondary":
        bg = palette["secondary"]
        text_color = palette["background"]
        border = palette["secondary"]
    elif req.variant == "outline":
        bg = "transparent"
        text_color = palette["primary"]
        border = palette["primary"]
    elif req.variant == "ghost":
        bg = "transparent"
        text_color = palette["primary"]
        border = "transparent"
    else:
        bg = palette["primary"]
        text_color = palette["background"]
        border = palette["primary"]

    # Component generators
    if req.component_type == "button":
        return f'''<button class="{size['padding']} {size['text']} font-medium rounded-lg transition-all hover:opacity-90" style="background-color: {bg}; color: {text_color}; border: 2px solid {border};">
    {content.get("text", "Button")}
</button>'''

    elif req.component_type == "card":
        return f'''<div class="rounded-xl p-6 shadow-lg" style="background-color: {palette['background']};">
    {f'<img src="{content.get("image")}" alt="" class="w-full h-48 object-cover rounded-lg mb-4">' if content.get("image") else ''}
    <h3 class="text-xl font-bold mb-2" style="color: {palette['text']};">{content.get("title", "Card Title")}</h3>
    <p class="opacity-70" style="color: {palette['text']};">{content.get("description", "Card description goes here.")}</p>
    {f'<a href="{content.get("cta_url", "#")}" class="inline-block mt-4 font-medium" style="color: {palette["secondary"]};">{content.get("cta_text", "Learn More")} →</a>' if content.get("cta_text") else ''}
</div>'''

    elif req.component_type == "navbar":
        links = content.get("links", [{"text": "Home", "url": "#"}, {"text": "About", "url": "#"}, {"text": "Contact", "url": "#"}])
        links_html = " ".join([f'<a href="{l.get("url", "#")}" class="hover:opacity-70 transition-opacity">{l.get("text", "")}</a>' for l in links])
        return f'''<nav class="flex items-center justify-between py-4 px-6" style="background-color: {palette['primary']}; color: {palette['background']};">
    <div class="text-xl font-bold">{content.get("logo_text", "Logo")}</div>
    <div class="hidden md:flex space-x-8">{links_html}</div>
    <button class="px-4 py-2 rounded-lg" style="background-color: {palette['secondary']}; color: {palette['background']};">{content.get("cta_text", "Get Started")}</button>
</nav>'''

    elif req.component_type == "input":
        return f'''<div class="space-y-2">
    <label class="block text-sm font-medium" style="color: {palette['text']};">{content.get("label", "Label")}</label>
    <input type="{content.get("type", "text")}" placeholder="{content.get("placeholder", "Enter text...")}" class="w-full {size['padding']} rounded-lg border focus:outline-none focus:ring-2" style="border-color: {palette['accent']}33; --tw-ring-color: {palette['secondary']};">
    {f'<p class="text-sm opacity-60">{content.get("helper")}</p>' if content.get("helper") else ''}
</div>'''

    elif req.component_type == "alert":
        alert_colors = {
            "info": {"bg": palette["primary"] + "20", "border": palette["primary"], "text": palette["primary"]},
            "success": {"bg": palette.get("success", "#28A745") + "20", "border": palette.get("success", "#28A745"), "text": palette.get("success", "#28A745")},
            "warning": {"bg": palette.get("warning", "#FFC107") + "20", "border": palette.get("warning", "#FFC107"), "text": palette.get("warning", "#FFC107")},
            "error": {"bg": palette.get("error", "#DC3545") + "20", "border": palette.get("error", "#DC3545"), "text": palette.get("error", "#DC3545")},
        }
        alert_style = alert_colors.get(content.get("type", "info"), alert_colors["info"])
        return f'''<div class="p-4 rounded-lg border-l-4" style="background-color: {alert_style['bg']}; border-color: {alert_style['border']};">
    <p class="font-medium" style="color: {alert_style['text']};">{content.get("title", "Alert Title")}</p>
    <p class="text-sm mt-1" style="color: {alert_style['text']}; opacity: 0.8;">{content.get("message", "Alert message goes here.")}</p>
</div>'''

    elif req.component_type == "badge":
        return f'''<span class="inline-flex items-center {size['padding']} {size['text']} font-medium rounded-full" style="background-color: {bg}; color: {text_color};">
    {content.get("text", "Badge")}
</span>'''

    elif req.component_type == "pricing_card":
        features = content.get("features", ["Feature 1", "Feature 2", "Feature 3"])
        features_html = "".join([f'<li class="flex items-center py-2"><svg class="w-5 h-5 mr-2" style="color: {palette["secondary"]};" fill="currentColor" viewBox="0 0 20 20"><path fill-rule="evenodd" d="M16.707 5.293a1 1 0 010 1.414l-8 8a1 1 0 01-1.414 0l-4-4a1 1 0 011.414-1.414L8 12.586l7.293-7.293a1 1 0 011.414 0z" clip-rule="evenodd"></path></svg>{f}</li>' for f in features])
        is_featured = content.get("featured", False)
        return f'''<div class="rounded-2xl p-8 {'border-2 scale-105' if is_featured else 'border'}" style="background-color: {palette['background']}; {'border-color: ' + palette['secondary'] + ';' if is_featured else 'border-color: ' + palette['accent'] + '33;'}">
    {f'<span class="inline-block px-3 py-1 text-sm font-semibold rounded-full mb-4" style="background-color: {palette["secondary"]}; color: {palette["background"]};">Most Popular</span>' if is_featured else ''}
    <h3 class="text-2xl font-bold" style="color: {palette['text']};">{content.get("name", "Plan")}</h3>
    <p class="opacity-60 mt-2" style="color: {palette['text']};">{content.get("description", "")}</p>
    <div class="my-6">
        <span class="text-4xl font-bold" style="color: {palette['text']};">{content.get("price", "$0")}</span>
        <span class="opacity-60">/{content.get("period", "month")}</span>
    </div>
    <ul class="space-y-1 mb-8" style="color: {palette['text']};">{features_html}</ul>
    <button class="w-full py-3 rounded-lg font-medium transition-all hover:opacity-90" style="background-color: {palette['secondary'] if is_featured else palette['primary']}; color: {palette['background']};">{content.get("cta_text", "Get Started")}</button>
</div>'''

    elif req.component_type == "testimonial_card":
        return f'''<div class="rounded-xl p-6 shadow-lg" style="background-color: {palette['background']};">
    <div class="flex items-center mb-4">
        <div class="w-12 h-12 rounded-full flex items-center justify-center text-lg font-bold" style="background-color: {palette['primary']}; color: {palette['background']};">{content.get("name", "A")[0]}</div>
        <div class="ml-4">
            <p class="font-semibold" style="color: {palette['text']};">{content.get("name", "Customer Name")}</p>
            <p class="text-sm opacity-60" style="color: {palette['text']};">{content.get("role", "Role")}</p>
        </div>
    </div>
    <p class="italic opacity-80" style="color: {palette['text']};">"{content.get("quote", "Great product!")}"</p>
    <div class="flex mt-4">{"⭐" * content.get("rating", 5)}</div>
</div>'''

    elif req.component_type == "feature_card":
        return f'''<div class="rounded-xl p-6 shadow-lg" style="background-color: {palette['background']};">
    <div class="w-12 h-12 mb-4 flex items-center justify-center rounded-lg" style="background-color: {palette['secondary']}20;">
        <span class="text-2xl">{content.get("icon", "✨")}</span>
    </div>
    <h3 class="text-xl font-bold mb-2" style="color: {palette['text']};">{content.get("title", "Feature")}</h3>
    <p class="opacity-70" style="color: {palette['text']};">{content.get("description", "Feature description.")}</p>
</div>'''

    elif req.component_type == "hero":
        return f'''<section class="py-20" style="background: linear-gradient(135deg, {palette['primary']}, {palette['accent']});">
    <div class="max-w-7xl mx-auto px-4 text-center">
        <h1 class="text-5xl font-bold mb-6" style="color: {palette['background']};">{content.get("title", "Welcome")}</h1>
        <p class="text-xl mb-8 opacity-90" style="color: {palette['background']};">{content.get("subtitle", "Your tagline here")}</p>
        <button class="px-8 py-4 rounded-lg font-semibold shadow-lg" style="background-color: {palette['secondary']}; color: {palette['background']};">{content.get("cta_text", "Get Started")}</button>
    </div>
</section>'''

    elif req.component_type == "cta_banner":
        return f'''<div class="rounded-xl p-8 text-center" style="background: linear-gradient(135deg, {palette['primary']}, {palette['accent']});">
    <h2 class="text-3xl font-bold mb-4" style="color: {palette['background']};">{content.get("title", "Ready to get started?")}</h2>
    <p class="text-lg mb-6 opacity-90" style="color: {palette['background']};">{content.get("subtitle", "Join thousands of happy customers.")}</p>
    <button class="px-8 py-3 rounded-lg font-semibold" style="background-color: {palette['background']}; color: {palette['primary']};">{content.get("cta_text", "Start Now")}</button>
</div>'''

    elif req.component_type == "form":
        fields = content.get("fields", [
            {"name": "name", "label": "Name", "type": "text"},
            {"name": "email", "label": "Email", "type": "email"},
            {"name": "message", "label": "Message", "type": "textarea"}
        ])
        fields_html = ""
        for field in fields:
            if field.get("type") == "textarea":
                fields_html += f'''<div class="space-y-2">
    <label class="block text-sm font-medium" style="color: {palette['text']};">{field.get("label", "")}</label>
    <textarea name="{field.get('name', '')}" rows="4" class="w-full px-4 py-3 rounded-lg border focus:outline-none focus:ring-2" style="border-color: {palette['accent']}33;"></textarea>
</div>'''
            else:
                fields_html += f'''<div class="space-y-2">
    <label class="block text-sm font-medium" style="color: {palette['text']};">{field.get("label", "")}</label>
    <input type="{field.get('type', 'text')}" name="{field.get('name', '')}" class="w-full px-4 py-3 rounded-lg border focus:outline-none focus:ring-2" style="border-color: {palette['accent']}33;">
</div>'''

        return f'''<form class="rounded-xl p-6 space-y-4" style="background-color: {palette['light']};">
    {fields_html}
    <button type="submit" class="w-full py-3 rounded-lg font-medium" style="background-color: {palette['primary']}; color: {palette['background']};">{content.get("submit_text", "Submit")}</button>
</form>'''

    elif req.component_type == "modal":
        return f'''<div class="fixed inset-0 flex items-center justify-center z-50" style="background-color: rgba(0,0,0,0.5);">
    <div class="rounded-xl p-6 max-w-md w-full mx-4 shadow-2xl" style="background-color: {palette['background']};">
        <div class="flex justify-between items-center mb-4">
            <h3 class="text-xl font-bold" style="color: {palette['text']};">{content.get("title", "Modal Title")}</h3>
            <button class="opacity-60 hover:opacity-100">&times;</button>
        </div>
        <p class="mb-6 opacity-70" style="color: {palette['text']};">{content.get("content", "Modal content goes here.")}</p>
        <div class="flex space-x-4">
            <button class="flex-1 py-2 rounded-lg" style="background-color: {palette['accent']}33; color: {palette['text']};">Cancel</button>
            <button class="flex-1 py-2 rounded-lg" style="background-color: {palette['primary']}; color: {palette['background']};">{content.get("cta_text", "Confirm")}</button>
        </div>
    </div>
</div>'''

    elif req.component_type == "tabs":
        tabs = content.get("tabs", [{"label": "Tab 1", "content": "Content 1"}, {"label": "Tab 2", "content": "Content 2"}])
        tabs_html = " ".join([f'<button class="px-4 py-2 font-medium {" border-b-2" if i == 0 else ""}" style="{'border-color: ' + palette['primary'] + '; color: ' + palette['primary'] + ';' if i == 0 else 'color: ' + palette['text'] + '; opacity: 0.6;'}">{t.get("label", "")}</button>' for i, t in enumerate(tabs)])
        return f'''<div>
    <div class="flex space-x-4 border-b" style="border-color: {palette['accent']}33;">
        {tabs_html}
    </div>
    <div class="py-4" style="color: {palette['text']};">
        {tabs[0].get("content", "") if tabs else ""}
    </div>
</div>'''

    elif req.component_type == "accordion":
        items = content.get("items", [{"title": "Section 1", "content": "Content 1"}])
        items_html = ""
        for item in items:
            items_html += f'''<details class="border-b" style="border-color: {palette['accent']}33;">
    <summary class="flex justify-between items-center cursor-pointer py-4">
        <span class="font-medium" style="color: {palette['text']};">{item.get("title", "")}</span>
        <span>▼</span>
    </summary>
    <div class="pb-4 opacity-70" style="color: {palette['text']};">{item.get("content", "")}</div>
</details>'''
        return f'''<div class="rounded-xl overflow-hidden border" style="border-color: {palette['accent']}33;">
    {items_html}
</div>'''

    elif req.component_type == "dropdown":
        options = content.get("options", ["Option 1", "Option 2", "Option 3"])
        options_html = "".join([f'<option value="{o}">{o}</option>' for o in options])
        return f'''<div class="space-y-2">
    <label class="block text-sm font-medium" style="color: {palette['text']};">{content.get("label", "Select")}</label>
    <select class="w-full px-4 py-3 rounded-lg border focus:outline-none focus:ring-2 appearance-none cursor-pointer" style="border-color: {palette['accent']}33; background-color: {palette['background']};">
        {options_html}
    </select>
</div>'''

    elif req.component_type == "footer":
        columns = content.get("columns", [])
        cols_html = ""
        for col in columns:
            links_html = "".join([f'<li><a href="{l.get("url", "#")}" class="hover:opacity-100 opacity-70">{l.get("text", "")}</a></li>' for l in col.get("links", [])])
            cols_html += f'''<div>
    <h4 class="font-semibold mb-4">{col.get("title", "")}</h4>
    <ul class="space-y-2">{links_html}</ul>
</div>'''
        return f'''<footer class="py-12" style="background-color: {palette['primary']}; color: {palette['background']};">
    <div class="max-w-7xl mx-auto px-4">
        <div class="grid grid-cols-2 md:grid-cols-4 gap-8 mb-8">{cols_html}</div>
        <div class="pt-8 border-t border-white/20 text-center opacity-70">
            <p>{content.get("copyright", "&copy; 2026 All rights reserved.")}</p>
        </div>
    </div>
</footer>'''

    return f"<!-- Component type '{req.component_type}' not implemented -->"

def generate_wireframe(req: WireframeRequest, palette: dict) -> str:
    """Generate a wireframe mockup"""
    # Wireframe uses grayscale with annotations
    wireframe_style = {
        "bg": "#FAFAFA",
        "element": "#E0E0E0",
        "text": "#666666",
        "annotation": "#3B82F6",
        "grid": "#F0F0F0"
    }

    # Section templates for different page types
    section_templates = {
        "navbar": '<div class="h-16 flex items-center justify-between px-6" style="background-color: {element};"><div class="w-24 h-8 rounded" style="background-color: #CCC;"></div><div class="flex space-x-4">{nav_items}</div><div class="w-20 h-8 rounded" style="background-color: #BBB;"></div></div>',
        "hero": '<div class="py-24 text-center" style="background-color: {element};"><div class="w-3/4 h-12 mx-auto mb-4 rounded" style="background-color: #CCC;"></div><div class="w-1/2 h-6 mx-auto mb-8 rounded" style="background-color: #DDD;"></div><div class="w-32 h-10 mx-auto rounded" style="background-color: #BBB;"></div></div>',
        "features": '<div class="py-16 px-6" style="background-color: {bg};"><div class="grid grid-cols-3 gap-8">{feature_items}</div></div>',
        "content": '<div class="py-12 px-6" style="background-color: {bg};"><div class="w-2/3 h-8 mb-4 rounded" style="background-color: #DDD;"></div><div class="space-y-2">{content_lines}</div></div>',
        "cta": '<div class="py-12 text-center" style="background-color: {element};"><div class="w-1/2 h-8 mx-auto mb-4 rounded" style="background-color: #CCC;"></div><div class="w-32 h-10 mx-auto rounded" style="background-color: #BBB;"></div></div>',
        "footer": '<div class="py-8 px-6" style="background-color: {element};"><div class="flex justify-between"><div class="w-24 h-6 rounded" style="background-color: #CCC;"></div><div class="flex space-x-4">{footer_items}</div></div></div>',
        "sidebar": '<div class="w-64 p-4" style="background-color: {element};"><div class="space-y-2">{sidebar_items}</div></div>',
        "grid": '<div class="py-12 px-6" style="background-color: {bg};"><div class="grid grid-cols-4 gap-4">{grid_items}</div></div>',
        "form": '<div class="py-12 px-6 max-w-md mx-auto" style="background-color: {bg};">{form_fields}<div class="w-full h-10 mt-4 rounded" style="background-color: #BBB;"></div></div>',
        "table": '<div class="py-12 px-6" style="background-color: {bg};"><div class="border rounded" style="border-color: #DDD;">{table_rows}</div></div>',
        "cards": '<div class="py-12 px-6" style="background-color: {bg};"><div class="grid grid-cols-3 gap-6">{card_items}</div></div>',
    }

    # Generate HTML elements
    nav_items = '<div class="w-16 h-4 rounded" style="background-color: #CCC;"></div>' * 4
    feature_items = ('<div class="p-6 rounded-lg" style="background-color: {element};">'
                    '<div class="w-12 h-12 mb-4 rounded" style="background-color: #CCC;"></div>'
                    '<div class="w-full h-4 mb-2 rounded" style="background-color: #DDD;"></div>'
                    '<div class="w-3/4 h-3 rounded" style="background-color: #EEE;"></div>'
                    '</div>').format(element=wireframe_style["element"]) * 3
    content_lines = '<div class="w-full h-3 rounded" style="background-color: #EEE;"></div>' * 5
    footer_items = '<div class="w-16 h-4 rounded" style="background-color: #CCC;"></div>' * 4
    sidebar_items = '<div class="w-full h-6 rounded" style="background-color: #CCC;"></div>' * 6
    grid_items = '<div class="aspect-square rounded" style="background-color: {element};"></div>'.format(element=wireframe_style["element"]) * 8
    form_fields = ('<div class="mb-4"><div class="w-20 h-3 mb-2 rounded" style="background-color: #CCC;"></div>'
                   '<div class="w-full h-10 rounded" style="background-color: #EEE;"></div></div>') * 3
    table_rows = ('<div class="flex border-b p-4" style="border-color: #EEE;">'
                  '<div class="flex-1 h-4 rounded" style="background-color: #DDD;"></div>'
                  '<div class="flex-1 h-4 rounded mx-4" style="background-color: #EEE;"></div>'
                  '<div class="flex-1 h-4 rounded" style="background-color: #EEE;"></div>'
                  '</div>') * 5
    card_items = ('<div class="rounded-lg overflow-hidden" style="background-color: {element};">'
                  '<div class="w-full h-32" style="background-color: #DDD;"></div>'
                  '<div class="p-4"><div class="w-3/4 h-4 mb-2 rounded" style="background-color: #CCC;"></div>'
                  '<div class="w-full h-3 rounded" style="background-color: #EEE;"></div></div>'
                  '</div>').format(element=wireframe_style["element"]) * 3

    html = f'''<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Wireframe - {req.page_type}</title>
    <script src="https://cdn.tailwindcss.com"></script>
    <style>
        .annotation {{ position: relative; }}
        .annotation::after {{
            content: attr(data-annotation);
            position: absolute;
            top: -24px;
            left: 0;
            font-size: 12px;
            color: {wireframe_style["annotation"]};
            background: white;
            padding: 2px 6px;
            border-radius: 4px;
            border: 1px solid {wireframe_style["annotation"]};
        }}
        {"" if not req.include_grid else ".grid-overlay { background-image: linear-gradient(to right, " + wireframe_style["grid"] + " 1px, transparent 1px), linear-gradient(to bottom, " + wireframe_style["grid"] + " 1px, transparent 1px); background-size: 20px 20px; }"}
    </style>
</head>
<body class="{"grid-overlay" if req.include_grid else ""}" style="background-color: {wireframe_style['bg']};">
    <div class="max-w-6xl mx-auto py-8 px-4">
        <div class="mb-4 text-center">
            <span class="inline-block px-3 py-1 rounded-full text-sm" style="background-color: {wireframe_style['annotation']}; color: white;">
                Wireframe: {req.page_type.title()} Page
            </span>
        </div>
'''

    for section_name in req.sections:
        section_name_lower = section_name.lower().replace(" ", "_")
        template = section_templates.get(section_name_lower, section_templates.get("content"))

        section_html = template.format(
            element=wireframe_style["element"],
            bg=wireframe_style["bg"],
            nav_items=nav_items,
            feature_items=feature_items,
            content_lines=content_lines,
            footer_items=footer_items,
            sidebar_items=sidebar_items,
            grid_items=grid_items,
            form_fields=form_fields,
            table_rows=table_rows,
            card_items=card_items
        )

        annotation_attr = f'data-annotation="{section_name}"' if req.annotations else ''
        html += f'''
        <div class="mb-8 {'annotation' if req.annotations else ''}" {annotation_attr}>
            {section_html}
        </div>
'''

    html += '''
    </div>
</body>
</html>'''

    return html

def generate_website_template(req: WebsiteTemplateRequest, palette: dict) -> dict:
    """Generate a complete website template with multiple pages"""
    style_config = WEB_STYLES.get(req.style, WEB_STYLES["modern"])

    if req.custom_colors:
        palette = {**palette, **req.custom_colors}

    # Default breakpoints
    breakpoints = req.responsive_breakpoints or {
        "sm": 640,
        "md": 768,
        "lg": 1024,
        "xl": 1280
    }

    # Page templates based on template type
    page_configs = {
        "landing": {
            "home": [
                ContentSection(type="hero", title="Transform Your Business", subtitle="The all-in-one solution for modern teams", cta_text="Get Started", background="gradient"),
                ContentSection(type="features", title="Why Choose Us", subtitle="Everything you need to succeed", items=[
                    {"icon": "🚀", "title": "Fast", "description": "Lightning fast performance"},
                    {"icon": "🔒", "title": "Secure", "description": "Enterprise-grade security"},
                    {"icon": "🎨", "title": "Beautiful", "description": "Stunning design out of the box"}
                ]),
                ContentSection(type="cta", title="Ready to get started?", subtitle="Join thousands of happy customers", cta_text="Start Free Trial"),
                ContentSection(type="footer", items=[
                    {"title": "Product", "links": [{"text": "Features", "url": "#"}, {"text": "Pricing", "url": "#"}]},
                    {"title": "Company", "links": [{"text": "About", "url": "#"}, {"text": "Blog", "url": "#"}]},
                ])
            ]
        },
        "saas": {
            "home": [
                ContentSection(type="hero", title="Supercharge Your Workflow", subtitle="AI-powered automation for the modern workplace", cta_text="Start Free", background="gradient"),
                ContentSection(type="stats", title="Trusted by thousands", items=[
                    {"value": "10K+", "label": "Active Users"},
                    {"value": "99.9%", "label": "Uptime"},
                    {"value": "24/7", "label": "Support"},
                    {"value": "50+", "label": "Integrations"}
                ]),
                ContentSection(type="features", title="Powerful Features", subtitle="Everything you need in one platform", items=[
                    {"icon": "⚡", "title": "Automation", "description": "Automate repetitive tasks"},
                    {"icon": "📊", "title": "Analytics", "description": "Deep insights into your data"},
                    {"icon": "🔗", "title": "Integrations", "description": "Connect with your favorite tools"}
                ]),
                ContentSection(type="pricing", title="Simple Pricing", subtitle="No hidden fees", items=[
                    {"name": "Starter", "price": "$9", "period": "month", "features": ["5 users", "10GB storage", "Basic support"], "cta_text": "Start Free"},
                    {"name": "Pro", "price": "$29", "period": "month", "features": ["25 users", "100GB storage", "Priority support", "API access"], "cta_text": "Get Started", "featured": True},
                    {"name": "Enterprise", "price": "$99", "period": "month", "features": ["Unlimited users", "Unlimited storage", "24/7 support", "Custom integrations"], "cta_text": "Contact Sales"}
                ]),
                ContentSection(type="testimonials", title="What Our Customers Say", items=[
                    {"name": "Sarah J.", "role": "CEO", "quote": "This tool transformed our workflow completely."},
                    {"name": "Mike R.", "role": "CTO", "quote": "Best investment we've made for our team."}
                ]),
                ContentSection(type="faq", title="Frequently Asked Questions", items=[
                    {"question": "How do I get started?", "answer": "Simply sign up for a free trial and follow our onboarding guide."},
                    {"question": "Can I cancel anytime?", "answer": "Yes, you can cancel your subscription at any time with no penalties."}
                ]),
                ContentSection(type="cta", title="Start Your Free Trial", subtitle="No credit card required", cta_text="Get Started Now"),
                ContentSection(type="footer", items=[
                    {"title": "Product", "links": [{"text": "Features", "url": "#"}, {"text": "Pricing", "url": "#"}, {"text": "Changelog", "url": "#"}]},
                    {"title": "Company", "links": [{"text": "About", "url": "#"}, {"text": "Blog", "url": "#"}, {"text": "Careers", "url": "#"}]},
                    {"title": "Legal", "links": [{"text": "Privacy", "url": "#"}, {"text": "Terms", "url": "#"}]}
                ])
            ]
        },
        "agency": {
            "home": [
                ContentSection(type="hero", title="We Build Digital Experiences", subtitle="Award-winning design and development agency", cta_text="View Our Work", background="gradient"),
                ContentSection(type="features", title="Our Services", items=[
                    {"icon": "🎨", "title": "Design", "description": "Beautiful, user-centered design"},
                    {"icon": "💻", "title": "Development", "description": "Robust, scalable solutions"},
                    {"icon": "📱", "title": "Mobile", "description": "Native and cross-platform apps"}
                ]),
                ContentSection(type="testimonials", title="Client Love", items=[
                    {"name": "Fortune 500 Company", "role": "Tech Director", "quote": "They delivered beyond our expectations."}
                ]),
                ContentSection(type="contact", title="Let's Work Together", subtitle="Tell us about your project"),
                ContentSection(type="footer", items=[
                    {"title": "Services", "links": [{"text": "Design", "url": "#"}, {"text": "Development", "url": "#"}]},
                    {"title": "Contact", "links": [{"text": "hello@agency.com", "url": "#"}]}
                ])
            ]
        }
    }

    # Get template config or use landing as default
    template_config = page_configs.get(req.template_type, page_configs["landing"])

    # Generate pages
    pages = {}
    for page_name in req.pages:
        if page_name in template_config:
            sections = template_config[page_name]
        else:
            # Default page structure
            sections = [
                ContentSection(type="hero", title=page_name.title(), subtitle=f"Welcome to the {page_name} page"),
                ContentSection(type="content", content="Page content goes here."),
                ContentSection(type="footer", items=[])
            ]

        page_req = LandingPageRequest(
            title=f"{req.template_type.title()} - {page_name.title()}",
            sections=sections,
            brand_id=req.brand_id,
            style=req.style,
            use_tailwind=req.use_tailwind,
            include_animations=True,
            custom_colors=req.custom_colors
        )

        pages[page_name] = generate_landing_page(page_req, palette)

    return pages


# =============================================================================
# API ENDPOINTS
# =============================================================================

@app.get("/health")
async def health():
    """Health check endpoint"""
    time_ctx = get_time_context()
    return {
        "status": "healthy",
        "agent": "THALIA",
        "role": "Creative Fleet Designer",
        "version": "1.1.0",
        "port": 8032,
        "fleet": "Creative Fleet",
        "capabilities": [
            "presentation_design",
            "document_formatting",
            "chart_creation",
            "thumbnail_design",
            "brand_template_application",
            "layout_optimization",
            "landing_page_design",
            "ui_component_design",
            "wireframe_mockups",
            "website_templates"
        ],
        "current_time": time_ctx["current_datetime"],
        "days_to_launch": time_ctx["days_to_launch"]
    }

@app.get("/time")
async def get_time():
    """Get current time context"""
    return get_time_context()

@app.get("/design-system")
async def get_design_system():
    """Return the design system configuration"""
    return {
        "color_palettes": COLOR_PALETTES,
        "typography": TYPOGRAPHY,
        "slide_layouts": SLIDE_LAYOUTS,
        "default_brand": "leveredge"
    }

@app.get("/palettes")
async def get_palettes():
    """List available color palettes"""
    return {
        "palettes": list(COLOR_PALETTES.keys()),
        "default": "leveredge"
    }

@app.get("/palette/{brand_id}")
async def get_palette_details(brand_id: str):
    """Get details of a specific palette"""
    if brand_id not in COLOR_PALETTES:
        raise HTTPException(status_code=404, detail=f"Palette '{brand_id}' not found")
    return {
        "brand_id": brand_id,
        "colors": COLOR_PALETTES[brand_id]
    }

@app.post("/design/presentation")
async def create_presentation(req: PresentationRequest):
    """Create a presentation from content"""
    time_ctx = get_time_context()

    try:
        palette = get_palette(req.brand_id)
        slides = req.content.get("slides", [])

        if not slides:
            raise HTTPException(status_code=400, detail="No slides provided in content")

        # Generate presentation
        file_path = generate_presentation(req.content, palette)

        # Generate preview URL (in production, this would be an actual accessible URL)
        preview_url = f"/files/{Path(file_path).name}"

        await notify_event_bus("presentation_created", {
            "slide_count": len(slides),
            "brand_id": req.brand_id,
            "file_path": file_path
        })

        return {
            "file_path": file_path,
            "preview_url": preview_url,
            "slide_count": len(slides),
            "brand_id": req.brand_id,
            "style": req.style,
            "agent": "THALIA",
            "timestamp": time_ctx["current_datetime"]
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Presentation creation failed: {str(e)}")

@app.post("/design/document")
async def format_document(req: DocumentRequest):
    """Format a document with brand styling"""
    time_ctx = get_time_context()

    try:
        palette = get_palette(req.brand_id)

        # Generate document
        file_path = generate_document(
            content=req.content,
            format=req.format,
            palette=palette,
            title=req.title,
            author=req.author,
            headers=req.headers
        )

        # Generate preview URL
        preview_url = f"/files/{Path(file_path).name}"

        await notify_event_bus("document_created", {
            "format": req.format,
            "brand_id": req.brand_id,
            "file_path": file_path
        })

        return {
            "file_path": file_path,
            "preview_url": preview_url,
            "format": req.format,
            "brand_id": req.brand_id,
            "agent": "THALIA",
            "timestamp": time_ctx["current_datetime"]
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Document creation failed: {str(e)}")

@app.post("/design/chart")
async def generate_chart_endpoint(req: ChartRequest):
    """Generate a chart/graph"""
    time_ctx = get_time_context()

    try:
        result = generate_chart(
            chart_type=req.type,
            data=req.data,
            style=req.style,
            title=req.title,
            width=req.width,
            height=req.height
        )

        await notify_event_bus("chart_created", {
            "chart_type": req.type,
            "style": req.style
        })

        return {
            "image_path": result["image_path"],
            "svg_path": result["svg_path"],
            "animation_path": result.get("animation_path"),
            "chart_type": req.type,
            "style": req.style,
            "agent": "THALIA",
            "timestamp": time_ctx["current_datetime"]
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Chart creation failed: {str(e)}")

@app.post("/design/thumbnail")
async def create_thumbnail(req: ThumbnailRequest):
    """Create a thumbnail image"""
    time_ctx = get_time_context()

    try:
        palette = get_palette(req.brand_id)

        sizes = generate_thumbnail(
            title=req.title,
            style=req.style,
            palette=palette,
            background_image=req.background_image,
            subtitle=req.subtitle
        )

        # Main image path (large size)
        image_path = sizes["large"]

        await notify_event_bus("thumbnail_created", {
            "title": req.title[:50],
            "style": req.style,
            "brand_id": req.brand_id
        })

        return {
            "image_path": image_path,
            "sizes": sizes,
            "style": req.style,
            "brand_id": req.brand_id,
            "agent": "THALIA",
            "timestamp": time_ctx["current_datetime"]
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Thumbnail creation failed: {str(e)}")

@app.post("/design/layout")
async def optimize_layout(req: LayoutRequest):
    """Optimize content layout (returns layout suggestions)"""
    time_ctx = get_time_context()

    palette = get_palette(req.brand_id)

    # Simple layout recommendations based on content type
    layouts = {
        "email": {
            "max_width": "600px",
            "padding": "20px",
            "font_size": "16px",
            "line_height": "1.6",
            "sections": ["header", "hero", "body", "cta", "footer"],
            "recommendations": [
                "Keep subject line under 50 characters",
                "Use single-column layout for mobile compatibility",
                "Place CTA above the fold",
                "Limit images to reduce load time"
            ]
        },
        "document": {
            "margins": "1 inch",
            "font_size": "12pt",
            "line_height": "1.5",
            "heading_hierarchy": ["h1: 24pt", "h2: 18pt", "h3: 14pt"],
            "recommendations": [
                "Use consistent heading styles",
                "Add page numbers for documents over 3 pages",
                "Include table of contents for long documents",
                "Use bullet points for lists of 3+ items"
            ]
        },
        "slide": {
            "aspect_ratio": "16:9",
            "max_bullets": 5,
            "font_min": "24pt",
            "recommendations": [
                "One idea per slide",
                "Maximum 6 words per bullet",
                "Use high-contrast colors",
                "Images should fill at least 40% of slide"
            ]
        },
        "social": {
            "optimal_sizes": {
                "instagram_square": "1080x1080",
                "instagram_story": "1080x1920",
                "linkedin": "1200x627",
                "twitter": "1200x675"
            },
            "recommendations": [
                "Use bold, readable text",
                "Keep text to 20% of image area",
                "Include brand colors consistently",
                "Test on mobile before posting"
            ]
        }
    }

    layout_config = layouts.get(req.content_type, layouts["document"])

    return {
        "content_type": req.content_type,
        "layout_config": layout_config,
        "brand_colors": palette,
        "typography": TYPOGRAPHY,
        "agent": "THALIA",
        "timestamp": time_ctx["current_datetime"]
    }

@app.get("/files/{filename}")
async def serve_file(filename: str):
    """Serve generated files"""
    file_path = f"{OUTPUT_DIR}/{filename}"
    if not Path(file_path).exists():
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(file_path)

@app.get("/templates")
async def list_templates():
    """List available templates"""
    return {
        "presentation_templates": [
            {"id": "pitch_deck", "name": "Pitch Deck", "slides": 10},
            {"id": "company_overview", "name": "Company Overview", "slides": 8},
            {"id": "project_update", "name": "Project Update", "slides": 6},
            {"id": "training", "name": "Training Deck", "slides": 12}
        ],
        "document_templates": [
            {"id": "proposal", "name": "Business Proposal"},
            {"id": "report", "name": "Status Report"},
            {"id": "memo", "name": "Internal Memo"},
            {"id": "brief", "name": "Project Brief"}
        ],
        "thumbnail_styles": ["gradient", "solid", "pattern", "image"]
    }

@app.post("/batch/charts")
async def batch_create_charts(charts: List[ChartRequest]):
    """Create multiple charts in batch"""
    time_ctx = get_time_context()
    results = []

    for chart_req in charts:
        try:
            result = generate_chart(
                chart_type=chart_req.type,
                data=chart_req.data,
                style=chart_req.style,
                title=chart_req.title,
                width=chart_req.width,
                height=chart_req.height
            )
            results.append({
                "success": True,
                "chart_type": chart_req.type,
                **result
            })
        except Exception as e:
            results.append({
                "success": False,
                "chart_type": chart_req.type,
                "error": str(e)
            })

    await notify_event_bus("batch_charts_created", {
        "total": len(charts),
        "successful": sum(1 for r in results if r["success"])
    })

    return {
        "results": results,
        "total": len(charts),
        "successful": sum(1 for r in results if r["success"]),
        "agent": "THALIA",
        "timestamp": time_ctx["current_datetime"]
    }

# =============================================================================
# WEB DESIGN ENDPOINTS
# =============================================================================

@app.post("/design/landing-page")
async def create_landing_page(request: LandingPageRequest):
    """
    Generate a complete landing page with specified sections.

    Sections can include: hero, features, cta, testimonials, pricing,
    faq, footer, stats, team, contact, custom.

    Returns HTML file with embedded Tailwind CSS.
    """
    time_ctx = get_time_context()

    # Get brand palette
    palette = get_brand_palette(request.brand_id)

    # Apply custom colors if provided
    if request.custom_colors:
        palette.update(request.custom_colors)

    try:
        # Generate the landing page HTML
        html_content = generate_landing_page(request, palette)

        # Save to file
        filename = f"landing_page_{uuid.uuid4().hex[:8]}.html"
        filepath = os.path.join(OUTPUT_DIR, filename)

        with open(filepath, "w", encoding="utf-8") as f:
            f.write(html_content)

        # Log cost (minimal for HTML generation)
        await cost_tracker.log_usage(
            operation="landing_page_generation",
            provider="internal",
            units=len(request.sections),
            unit_type="sections",
            cost_usd=0.0
        )

        # Notify event bus
        await notify_event_bus("landing_page_created", {
            "title": request.title,
            "sections": len(request.sections),
            "style": request.style,
            "filename": filename
        })

        return {
            "success": True,
            "file_path": filepath,
            "filename": filename,
            "preview_url": f"/files/{filename}",
            "sections_generated": len(request.sections),
            "style": request.style,
            "tailwind_enabled": request.use_tailwind,
            "responsive": request.responsive,
            "agent": "THALIA",
            "timestamp": time_ctx["current_datetime"]
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Landing page generation failed: {str(e)}")


@app.post("/design/ui-component")
async def create_ui_component(request: UIComponentRequest):
    """
    Generate individual UI components with Tailwind CSS.

    Component types: button, card, form, navbar, footer, modal, alert,
    badge, input, dropdown, tabs, accordion, pricing_card, testimonial_card,
    feature_card, hero, cta_banner.

    Variants: primary, secondary, outline, ghost (depends on component)
    Sizes: sm, md, lg, xl
    """
    time_ctx = get_time_context()

    # Get brand palette
    palette = get_brand_palette(request.brand_id)

    try:
        # Generate the component HTML
        html_content = generate_ui_component(request, palette)

        # Save to file
        filename = f"component_{request.component_type}_{uuid.uuid4().hex[:8]}.html"
        filepath = os.path.join(OUTPUT_DIR, filename)

        with open(filepath, "w", encoding="utf-8") as f:
            f.write(html_content)

        # Log cost
        await cost_tracker.log_usage(
            operation="ui_component_generation",
            provider="internal",
            units=1,
            unit_type="component",
            cost_usd=0.0
        )

        # Notify event bus
        await notify_event_bus("ui_component_created", {
            "component_type": request.component_type,
            "variant": request.variant,
            "size": request.size,
            "filename": filename
        })

        return {
            "success": True,
            "file_path": filepath,
            "filename": filename,
            "preview_url": f"/files/{filename}",
            "component_type": request.component_type,
            "variant": request.variant,
            "size": request.size,
            "tailwind_enabled": request.use_tailwind,
            "dark_mode": request.dark_mode,
            "agent": "THALIA",
            "timestamp": time_ctx["current_datetime"]
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"UI component generation failed: {str(e)}")


@app.post("/design/wireframe")
async def create_wireframe(request: WireframeRequest):
    """
    Generate wireframe/mockup for a page layout.

    Page types: landing, dashboard, blog, ecommerce, portfolio, saas, app.

    Outputs HTML wireframe with annotations showing section purposes.
    """
    time_ctx = get_time_context()

    # Get brand palette
    palette = get_brand_palette(request.brand_id)

    try:
        # Generate the wireframe
        html_content = generate_wireframe(request, palette)

        # Save to file
        filename = f"wireframe_{request.page_type}_{uuid.uuid4().hex[:8]}.html"
        filepath = os.path.join(OUTPUT_DIR, filename)

        with open(filepath, "w", encoding="utf-8") as f:
            f.write(html_content)

        # Log cost
        await cost_tracker.log_usage(
            operation="wireframe_generation",
            provider="internal",
            units=len(request.sections),
            unit_type="sections",
            cost_usd=0.0
        )

        # Notify event bus
        await notify_event_bus("wireframe_created", {
            "page_type": request.page_type,
            "sections": request.sections,
            "format": request.format,
            "filename": filename
        })

        return {
            "success": True,
            "file_path": filepath,
            "filename": filename,
            "preview_url": f"/files/{filename}",
            "page_type": request.page_type,
            "sections": request.sections,
            "format": request.format,
            "annotations_enabled": request.annotations,
            "grid_enabled": request.include_grid,
            "agent": "THALIA",
            "timestamp": time_ctx["current_datetime"]
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Wireframe generation failed: {str(e)}")


@app.post("/design/website-template")
async def create_website_template(request: WebsiteTemplateRequest):
    """
    Generate a complete website template with multiple pages.

    Template types: landing, portfolio, blog, business, saas, agency, startup.

    Returns a bundle of HTML files for each requested page.
    """
    time_ctx = get_time_context()

    # Get brand palette
    palette = get_brand_palette(request.brand_id)

    # Apply custom colors if provided
    if request.custom_colors:
        palette.update(request.custom_colors)

    try:
        # Generate the template (returns dict of page_name -> html_content)
        template_result = generate_website_template(request, palette)

        # Create output directory for template
        template_id = uuid.uuid4().hex[:8]
        template_dir = os.path.join(OUTPUT_DIR, f"template_{template_id}")
        Path(template_dir).mkdir(parents=True, exist_ok=True)

        # Save each page
        saved_files = []
        for page_name, html_content in template_result.get("pages", {}).items():
            filename = f"{page_name}.html"
            filepath = os.path.join(template_dir, filename)

            with open(filepath, "w", encoding="utf-8") as f:
                f.write(html_content)

            saved_files.append({
                "page": page_name,
                "filename": filename,
                "path": filepath,
                "preview_url": f"/files/template_{template_id}/{filename}"
            })

        # Log cost
        await cost_tracker.log_usage(
            operation="website_template_generation",
            provider="internal",
            units=len(request.pages or ["home"]),
            unit_type="pages",
            cost_usd=0.0
        )

        # Notify event bus
        await notify_event_bus("website_template_created", {
            "template_type": request.template_type,
            "pages": request.pages,
            "style": request.style,
            "template_id": template_id
        })

        return {
            "success": True,
            "template_id": template_id,
            "template_dir": template_dir,
            "pages": saved_files,
            "template_type": request.template_type,
            "style": request.style,
            "tailwind_enabled": request.use_tailwind,
            "javascript_enabled": request.include_javascript,
            "agent": "THALIA",
            "timestamp": time_ctx["current_datetime"]
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Website template generation failed: {str(e)}")


# =============================================================================
# MAIN
# =============================================================================

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8032)
